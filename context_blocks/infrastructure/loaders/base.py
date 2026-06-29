"""Base document loader interface."""

import hashlib
from dataclasses import dataclass, field
from pathlib import Path


@dataclass
class LoadedDocument:
    """A document loaded and ready for processing."""

    filename: str
    filepath: Path
    content: str
    file_type: str  # "pdf", "md", "txt"
    size_bytes: int


async def load_document(filepath: Path) -> LoadedDocument:
    """Load a single document from disk.

    Supports: .md, .txt, .pdf

    Args:
        filepath: Path to the document file.

    Returns:
        LoadedDocument with extracted text content.

    Raises:
        DocumentLoadError: If the file cannot be read or parsed.
    """
    from context_blocks.exceptions import DocumentLoadError

    if not filepath.exists():
        raise DocumentLoadError(f"File not found: {filepath}")

    suffix = filepath.suffix.lower()

    if suffix in (".md", ".txt"):
        try:
            content = filepath.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            try:
                content = filepath.read_text(encoding="latin-1")
            except Exception as e:
                raise DocumentLoadError(f"Cannot decode {filepath}: {e}") from e
    elif suffix == ".pdf":
        content = _extract_pdf_text(filepath)
    elif suffix == ".docx":
        content = _extract_docx_text(filepath)
    elif suffix == ".pptx":
        content = _extract_pptx_text(filepath)
    elif suffix in (".html", ".htm"):
        content = _extract_html_text(filepath)
    else:
        raise DocumentLoadError(f"Unsupported file type: {suffix}")

    # Skip empty documents
    content = content.strip()
    if not content:
        raise DocumentLoadError(f"Document is empty: {filepath}")

    return LoadedDocument(
        filename=filepath.name,
        filepath=filepath,
        content=content,
        file_type=suffix.lstrip("."),
        size_bytes=filepath.stat().st_size,
    )


def _extract_pdf_text(filepath: Path) -> str:
    """Extract text from a PDF file."""
    from context_blocks.exceptions import DocumentLoadError

    try:
        from pypdf import PdfReader

        reader = PdfReader(str(filepath))
        pages = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
        return "\n\n".join(pages)
    except Exception as e:
        raise DocumentLoadError(f"Failed to extract PDF text: {e}") from e


def _extract_docx_text(filepath: Path) -> str:
    """Extract text from a Word document."""
    from context_blocks.exceptions import DocumentLoadError

    try:
        from docx import Document

        doc = Document(str(filepath))
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n\n".join(parts)
    except ImportError:
        raise DocumentLoadError(
            "python-docx not installed. Install with: pip install 'context-blocks[docx]'"
        )
    except Exception as e:
        raise DocumentLoadError(f"Failed to extract DOCX text: {e}") from e


def _extract_pptx_text(filepath: Path) -> str:
    """Extract text from a PowerPoint presentation."""
    from context_blocks.exceptions import DocumentLoadError

    try:
        from pptx import Presentation

        prs = Presentation(str(filepath))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_texts.append(para.text.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            slide_texts.append(" | ".join(cells))
            if slide_texts:
                parts.append(f"## Slide {i}\n\n" + "\n\n".join(slide_texts))
        return "\n\n".join(parts)
    except ImportError:
        raise DocumentLoadError(
            "python-pptx not installed. Install with: pip install 'context-blocks[pptx]'"
        )
    except Exception as e:
        raise DocumentLoadError(f"Failed to extract PPTX text: {e}") from e


def _extract_html_text(filepath: Path) -> str:
    """Extract text from an HTML file (strips tags, nav, scripts)."""
    from context_blocks.exceptions import DocumentLoadError

    try:
        from html.parser import HTMLParser

        raw = filepath.read_text(encoding="utf-8")

        class _TextExtractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self.parts: list[str] = []
                self._skip = False
                self._skip_tags = {"script", "style", "nav", "footer", "header"}

            def handle_starttag(self, tag, attrs):
                if tag in self._skip_tags:
                    self._skip = True

            def handle_endtag(self, tag):
                if tag in self._skip_tags:
                    self._skip = False

            def handle_data(self, data):
                if not self._skip:
                    text = data.strip()
                    if text:
                        self.parts.append(text)

        extractor = _TextExtractor()
        extractor.feed(raw)
        return "\n\n".join(extractor.parts)
    except Exception as e:
        raise DocumentLoadError(f"Failed to extract HTML text: {e}") from e


async def load_documents_from_directory(directory: Path) -> list[LoadedDocument]:
    """Load all supported documents from a directory (recursive).

    Supports nested subdirectories — folder names are preserved as context signals.

    Args:
        directory: Path to directory containing documents.

    Returns:
        List of loaded documents, sorted by path.
    """
    from context_blocks.exceptions import DocumentLoadError

    import structlog

    logger = structlog.get_logger(__name__)

    if not directory.is_dir():
        raise DocumentLoadError(f"Not a directory: {directory}")

    supported_extensions = {".md", ".txt", ".pdf", ".docx", ".pptx", ".html", ".htm"}
    files = sorted(
        f for f in directory.rglob("*")
        if f.is_file()
        and f.suffix.lower() in supported_extensions
        and not f.name.startswith(".")
    )

    documents = []
    for filepath in files:
        try:
            doc = await load_document(filepath)
            documents.append(doc)
            logger.debug("document_loaded", filename=doc.filename, size=doc.size_bytes)
        except DocumentLoadError as e:
            logger.warning("document_load_failed", filename=filepath.name, error=str(e))
            continue

    # Dedup: detect documents with identical or near-identical content
    # (e.g., cost-calculation.txt and cost-calculation.md are the same doc in different formats)
    documents, skipped = _dedup_documents(documents)
    if skipped:
        for s in skipped:
            logger.info(
                "document_skipped_duplicate",
                filename=s["skipped"],
                duplicate_of=s["kept"],
            )

    logger.info("documents_loaded", count=len(documents), directory=str(directory))
    return documents


def _content_hash(text: str) -> str:
    """Hash normalized content for dedup comparison."""
    # Normalize: lowercase, collapse whitespace, strip formatting markers
    normalized = " ".join(text.lower().split())
    return hashlib.md5(normalized.encode()).hexdigest()


def _dedup_documents(documents: list[LoadedDocument]) -> tuple[list[LoadedDocument], list[dict]]:
    """Remove documents with duplicate content, keeping the first occurrence.

    Two docs are considered duplicates if their normalized content hash matches.
    When a duplicate is found, the larger file (more content) is kept.

    Returns:
        Tuple of (deduplicated docs, list of skipped info dicts).
    """
    seen_hashes: dict[str, LoadedDocument] = {}
    kept: list[LoadedDocument] = []
    skipped: list[dict] = []

    for doc in documents:
        h = _content_hash(doc.content)
        if h in seen_hashes:
            existing = seen_hashes[h]
            # Keep the larger one (more content)
            if doc.size_bytes > existing.size_bytes:
                # Replace: keep new, skip existing
                kept = [d for d in kept if d.filename != existing.filename]
                kept.append(doc)
                skipped.append({"skipped": existing.filename, "kept": doc.filename})
                seen_hashes[h] = doc
            else:
                skipped.append({"skipped": doc.filename, "kept": existing.filename})
        else:
            seen_hashes[h] = doc
            kept.append(doc)

    return kept, skipped
